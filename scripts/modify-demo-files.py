#!/usr/bin/env python3
"""Apply test modifications to all demo_files for change detection testing."""

import csv
import json
import os
import struct
import zlib

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEMO_DIR = os.path.join(PROJECT_ROOT, "demo_files")


def modify_text_files():
    with open(os.path.join(DEMO_DIR, "welcome.txt"), "w", encoding="utf-8") as f:
        f.write(
            "Welcome to the File Integrity Monitor demo folder.\n"
            "MODIFIED: This line was changed for testing.\n"
            "Line three — unchanged baseline content.\n"
            "NEW LINE: Added after baseline.\n"
        )

    with open(os.path.join(DEMO_DIR, "confidential.txt"), "w", encoding="utf-8") as f:
        f.write(
            "CONFIDENTIAL — Internal use only\n"
            "Project: File Integrity Monitoring System\n"
            "Owner: External Contractor (CHANGED)\n"
            "Status: COMPROMISED — unauthorized edit detected\n"
        )

    with open(os.path.join(DEMO_DIR, "readme.md"), "w", encoding="utf-8") as f:
        f.write(
            "# Demo Folder — MODIFIED\n\n"
            "This readme was edited to test change detection.\n\n"
            "## How to test\n\n"
            "1. Add this folder as a monitor in the app.\n"
            "2. Run **Check Now** — you should see this file listed.\n"
            "3. Use **What changed** and **Full file side by side** views.\n\n"
            "## Files included\n\n"
            "- All demo types were modified on purpose.\n"
        )

    with open(os.path.join(DEMO_DIR, "config.json"), "w", encoding="utf-8") as f:
        json.dump(
            {
                "app": "File Integrity Monitor",
                "version": "2.0.0",
                "monitoring": {"enabled": False, "interval_minutes": 5},
                "demo": False,
                "tampered": True,
            },
            f,
            indent=2,
        )
        f.write("\n")

    with open(os.path.join(DEMO_DIR, "data.csv"), "w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["id", "name", "role", "status"])
        writer.writerow(["1", "Alice", "Admin", "active"])
        writer.writerow(["2", "Bob", "Superuser", "suspended"])
        writer.writerow(["3", "Carol", "Auditor", "active"])
        writer.writerow(["4", "Dave", "Intruder", "unknown"])


def modify_rtf():
    content = r"""{\rtf1\ansi
{\b Demo RTF Document — MODIFIED}\par
This rich text file was tampered with for testing.\par
\par
Changed paragraph — not the baseline version.\par
Contact: attacker@example.com\par
}"""
    with open(os.path.join(DEMO_DIR, "memo.rtf"), "w", encoding="utf-8") as f:
        f.write(content)


def modify_docx():
    path = os.path.join(DEMO_DIR, "project-memo.docx")
    doc = Document(path)
    doc.paragraphs[1].text = (
        "MODIFIED: This Word document was edited after the baseline was created."
    )
    doc.paragraphs[2].text = "Approved by: Unknown Party (CHANGED)"
    table = doc.tables[0]
    table.cell(1, 1).text = "Tampered"
    doc.save(path)


def modify_xlsx():
    path = os.path.join(DEMO_DIR, "budget.xlsx")
    wb = load_workbook(path)
    ws = wb["Budget"]
    ws["B2"] = 99999
    ws["C3"] = 0
    summary = wb["Summary"]
    summary["A3"] = "Unauthorized Editor"
    wb.save(path)


def modify_pptx():
    path = os.path.join(DEMO_DIR, "overview.pptx")
    prs = Presentation(path)
    prs.slides[0].placeholders[1].text = (
        "MODIFIED slide text.\n"
        "This change should appear in the integrity check."
    )
    prs.slides[1].shapes.title.text = "Key Points — CHANGED"
    prs.save(path)


def modify_pdf():
    path = os.path.join(DEMO_DIR, "security-policy.pdf")
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, height - 72, "Security Policy — MODIFIED")
    c.setFont("Helvetica", 11)
    lines = [
        "This PDF was altered after baseline creation.",
        "",
        "Section 1: Access Control — DISABLED",
        "All users now have unrestricted access. (TAMPERED)",
        "",
        "Section 2: File Monitoring",
        "Monitoring interval changed to never.",
    ]
    y = height - 110
    for line in lines:
        c.drawString(72, y, line)
        y -= 16
    c.save()


def _png_chunk(tag: bytes, data: bytes) -> bytes:
    crc = zlib.crc32(tag + data) & 0xFFFFFFFF
    return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)


def modify_png():
    """Replace blue PNG with red PNG."""
    width, height = 48, 48
    raw = b""
    for _ in range(height):
        raw += b"\x00" + bytes([200, 40, 40]) * width
    compressed = zlib.compress(raw, 9)
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    png = b"\x89PNG\r\n\x1a\n"
    png += _png_chunk(b"IHDR", ihdr)
    png += _png_chunk(b"IDAT", compressed)
    png += _png_chunk(b"IEND", b"")
    with open(os.path.join(DEMO_DIR, "logo.png"), "wb") as f:
        f.write(png)


def modify_jpeg():
    """Different minimal JPEG (slightly altered bytes)."""
    jpeg = bytes([
        0xFF, 0xD8, 0xFF, 0xE0, 0x00, 0x10, 0x4A, 0x46, 0x49, 0x46, 0x00, 0x01,
        0x01, 0x00, 0x00, 0x01, 0x00, 0x01, 0x00, 0x00, 0xFF, 0xDB, 0x00, 0x43,
        0x00, 0x08, 0x06, 0x06, 0x07, 0x06, 0x05, 0x08, 0x07, 0x07, 0x07, 0x09,
        0x09, 0x08, 0x0A, 0x0C, 0x14, 0x0D, 0x0C, 0x0B, 0x0B, 0x0C, 0x19, 0x12,
        0x13, 0x0F, 0x14, 0x1D, 0x1A, 0x1F, 0x1E, 0x1D, 0x1A, 0x1C, 0x1C, 0x20,
        0x24, 0x2E, 0x27, 0x20, 0x22, 0x2C, 0x23, 0x1C, 0x1C, 0x28, 0x37, 0x29,
        0x2C, 0x30, 0x31, 0x34, 0x34, 0x34, 0x1F, 0x27, 0x39, 0x3D, 0x38, 0x32,
        0x3C, 0x2E, 0x33, 0x34, 0x32, 0xFF, 0xC0, 0x00, 0x11, 0x08, 0x00, 0x40,
        0x00, 0x40, 0x03, 0x01, 0x22, 0x00, 0x02, 0x11, 0x01, 0x03, 0x11, 0x01,
        0xFF, 0xC4, 0x00, 0x14, 0x00, 0x01, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00,
        0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x08, 0xFF, 0xC4,
        0x00, 0x14, 0x10, 0x01, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00,
        0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0xFF, 0xDA, 0x00, 0x0C,
        0x03, 0x01, 0x00, 0x02, 0x11, 0x03, 0x11, 0x00, 0x3F, 0x00, 0xAA, 0xFF,
        0xD9,
    ])
    with open(os.path.join(DEMO_DIR, "badge.jpg"), "wb") as f:
        f.write(jpeg)


def main():
    modify_text_files()
    modify_rtf()
    modify_docx()
    modify_xlsx()
    modify_pptx()
    modify_pdf()
    modify_png()
    modify_jpeg()
    print(f"Modified all demo files in {DEMO_DIR}")


if __name__ == "__main__":
    main()
