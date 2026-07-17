#!/usr/bin/env python3
"""Generate demo files in demo_files/ for testing all supported file types."""

import csv
import json
import os
import struct
import zlib

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEMO_DIR = os.path.join(PROJECT_ROOT, "demo_files")


def write_text_files():
    with open(os.path.join(DEMO_DIR, "welcome.txt"), "w", encoding="utf-8") as f:
        f.write(
            "Welcome to the File Integrity Monitor demo folder.\n"
            "This is the original welcome message.\n"
            "Line three — unchanged baseline content.\n"
        )

    with open(os.path.join(DEMO_DIR, "confidential.txt"), "w", encoding="utf-8") as f:
        f.write(
            "CONFIDENTIAL — Internal use only\n"
            "Project: File Integrity Monitoring System\n"
            "Owner: Security Team\n"
            "Status: Baseline snapshot active\n"
        )

    with open(os.path.join(DEMO_DIR, "readme.md"), "w", encoding="utf-8") as f:
        f.write(
            "# Demo Folder\n\n"
            "Use this folder to test monitoring.\n\n"
            "## How to test\n\n"
            "1. Add this folder as a monitor in the app.\n"
            "2. Run **Check Now** — everything should be clean.\n"
            "3. Edit any file below, then run **Check Now** again.\n"
            "4. Expand the changed file to see **What changed** or **Full file side by side**.\n\n"
            "## Files included\n\n"
            "- Text, Markdown, CSV, JSON\n"
            "- Word, Excel, PowerPoint\n"
            "- PDF and RTF\n"
            "- PNG and JPEG images\n"
        )

    with open(os.path.join(DEMO_DIR, "config.json"), "w", encoding="utf-8") as f:
        json.dump(
            {
                "app": "File Integrity Monitor",
                "version": "2.0.0",
                "monitoring": {"enabled": True, "interval_minutes": 20},
                "demo": True,
            },
            f,
            indent=2,
        )
        f.write("\n")

    with open(os.path.join(DEMO_DIR, "data.csv"), "w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["id", "name", "role", "status"])
        writer.writerow(["1", "Alice", "Admin", "active"])
        writer.writerow(["2", "Bob", "Analyst", "active"])
        writer.writerow(["3", "Carol", "Auditor", "active"])


def write_rtf():
    content = r"""{\rtf1\ansi
{\b Demo RTF Document}\par
This rich text file is monitored by the integrity system.\par
\par
Original paragraph — baseline version.\par
Contact: security@example.com\par
}"""
    with open(os.path.join(DEMO_DIR, "memo.rtf"), "w", encoding="utf-8") as f:
        f.write(content)


def write_docx():
    doc = Document()
    doc.add_heading("Project Memo", level=1)
    doc.add_paragraph(
        "This Word document is part of the demo baseline. "
        "Edit a paragraph after baseline creation to test the change viewer."
    )
    doc.add_paragraph("Approved by: Security Team")
    doc.add_paragraph("Effective date: Baseline snapshot")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Monitoring"
    table.cell(1, 1).text = "Active"
    doc.save(os.path.join(DEMO_DIR, "project-memo.docx"))


def write_xlsx():
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws["A1"] = "Category"
    ws["B1"] = "Q1"
    ws["C1"] = "Q2"
    rows = [
        ("Infrastructure", 12000, 13500),
        ("Licenses", 4500, 4500),
        ("Training", 2000, 2500),
    ]
    for row_index, row in enumerate(rows, start=2):
        for col_index, value in enumerate(row, start=1):
            ws.cell(row=row_index, column=col_index, value=value)

    summary = wb.create_sheet("Summary")
    summary["A1"] = "Total"
    summary["B1"] = "=Budget!B2+Budget!B3+Budget!B4"
    summary["A2"] = "Owner"
    summary["A3"] = "Finance Team"

    wb.save(os.path.join(DEMO_DIR, "budget.xlsx"))


def write_pptx():
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Integrity Monitor Demo"
    slide.placeholders[1].text = (
        "Baseline presentation slide.\n"
        "Modify this text to trigger a change alert."
    )

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Key Points"
    slide2.placeholders[1].text = "• All file types supported\n• Side-by-side diff\n• Restore from backup"

    prs.save(os.path.join(DEMO_DIR, "overview.pptx"))


def write_pdf():
    path = os.path.join(DEMO_DIR, "security-policy.pdf")
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter

    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, height - 72, "Security Policy — Demo")
    c.setFont("Helvetica", 11)
    lines = [
        "This PDF is included in the demo folder for integrity testing.",
        "",
        "Section 1: Access Control",
        "All users must use strong authentication.",
        "",
        "Section 2: File Monitoring",
        "Critical documents are checked every 20 minutes.",
        "",
        "Baseline version — edit this PDF or replace it to test change detection.",
    ]
    y = height - 110
    for line in lines:
        c.drawString(72, y, line)
        y -= 16

    c.showPage()
    c.setFont("Helvetica-Bold", 14)
    c.drawString(72, height - 72, "Appendix A")
    c.setFont("Helvetica", 11)
    c.drawString(72, height - 100, "Contact: security@example.com")
    c.save()


def _png_chunk(tag: bytes, data: bytes) -> bytes:
    crc = zlib.crc32(tag + data) & 0xFFFFFFFF
    return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)


def write_png():
    """Minimal valid 48x48 PNG (blue square) without Pillow."""
    width, height = 48, 48
    raw = b""
    for _ in range(height):
        raw += b"\x00" + bytes([30, 80, 160]) * width

    compressed = zlib.compress(raw, 9)
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    png = b"\x89PNG\r\n\x1a\n"
    png += _png_chunk(b"IHDR", ihdr)
    png += _png_chunk(b"IDAT", compressed)
    png += _png_chunk(b"IEND", b"")

    with open(os.path.join(DEMO_DIR, "logo.png"), "wb") as f:
        f.write(png)


def write_jpeg():
    """Minimal valid 32x32 JPEG (solid green) without Pillow."""
    # Tiny baseline JPEG
    jpeg = bytes([
        0xFF, 0xD8, 0xFF, 0xE0, 0x00, 0x10, 0x4A, 0x46, 0x49, 0x46, 0x00, 0x01,
        0x01, 0x00, 0x00, 0x01, 0x00, 0x01, 0x00, 0x00, 0xFF, 0xDB, 0x00, 0x43,
        0x00, 0x08, 0x06, 0x06, 0x07, 0x06, 0x05, 0x08, 0x07, 0x07, 0x07, 0x09,
        0x09, 0x08, 0x0A, 0x0C, 0x14, 0x0D, 0x0C, 0x0B, 0x0B, 0x0C, 0x19, 0x12,
        0x13, 0x0F, 0x14, 0x1D, 0x1A, 0x1F, 0x1E, 0x1D, 0x1A, 0x1C, 0x1C, 0x20,
        0x24, 0x2E, 0x27, 0x20, 0x22, 0x2C, 0x23, 0x1C, 0x1C, 0x28, 0x37, 0x29,
        0x2C, 0x30, 0x31, 0x34, 0x34, 0x34, 0x1F, 0x27, 0x39, 0x3D, 0x38, 0x32,
        0x3C, 0x2E, 0x33, 0x34, 0x32, 0xFF, 0xC0, 0x00, 0x11, 0x08, 0x00, 0x20,
        0x00, 0x20, 0x03, 0x01, 0x22, 0x00, 0x02, 0x11, 0x01, 0x03, 0x11, 0x01,
        0xFF, 0xC4, 0x00, 0x14, 0x00, 0x01, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00,
        0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x08, 0xFF, 0xC4,
        0x00, 0x14, 0x10, 0x01, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00,
        0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0x00, 0xFF, 0xDA, 0x00, 0x0C,
        0x03, 0x01, 0x00, 0x02, 0x11, 0x03, 0x11, 0x00, 0x3F, 0x00, 0xBF, 0x80,
        0xFF, 0xD9,
    ])
    with open(os.path.join(DEMO_DIR, "badge.jpg"), "wb") as f:
        f.write(jpeg)


def write_testing_guide():
    guide = """# Demo Files — Testing Guide

This folder contains sample files for every major supported type.

## Setup

1. In the app, add this folder as a monitor:
   `{demo_path}`
2. Wait for the baseline to finish (or run Check Now once).

## Quick tests

| File | How to trigger a change |
|------|-------------------------|
| welcome.txt | Edit line 2 |
| confidential.txt | Add a new line |
| readme.md | Change a heading |
| config.json | Change `"demo"` to `false` |
| data.csv | Change Bob's role |
| memo.rtf | Edit a paragraph |
| project-memo.docx | Change "Security Team" |
| budget.xlsx | Change a cell in Budget sheet |
| overview.pptx | Edit slide 1 subtitle |
| security-policy.pdf | Regenerate or edit text |
| logo.png | Replace with another image |
| badge.jpg | Replace with another image |

## What to look for

- **What changed** tab — summary cards + changed sections only
- **Full file side by side** — complete original vs current with highlights
- **Undo changes** — restore from baseline backup

## After testing restore

Click **Undo changes — restore original file** on any changed item, then run Check Now again to confirm it is clean.
""".format(demo_path=DEMO_DIR)

    with open(os.path.join(DEMO_DIR, "TESTING.md"), "w", encoding="utf-8") as f:
        f.write(guide)


def main():
    os.makedirs(DEMO_DIR, exist_ok=True)
    write_text_files()
    write_rtf()
    write_docx()
    write_xlsx()
    write_pptx()
    write_pdf()
    write_png()
    write_jpeg()
    write_testing_guide()
    print(f"Created demo files in {DEMO_DIR}")
    for name in sorted(os.listdir(DEMO_DIR)):
        path = os.path.join(DEMO_DIR, name)
        print(f"  {name} ({os.path.getsize(path)} bytes)")


if __name__ == "__main__":
    main()
