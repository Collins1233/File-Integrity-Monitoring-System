import difflib

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    import xlrd
except ImportError:
    xlrd = None


def read_text_file(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            return file.readlines()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="latin-1") as file:
                return file.readlines()
        except Exception:
            return []
    except Exception:
        return []


def compare_text(old_file, new_file):
    old_lines = read_text_file(old_file)
    new_lines = read_text_file(new_file)
    return list(
        difflib.unified_diff(
            old_lines,
            new_lines,
            fromfile="Baseline Version",
            tofile="Current Version",
            lineterm="",
        )
    )


def _format_line_comparison(before_lines, after_lines, preview_type, document_type=None, extra=None):
    before = [line.rstrip("\n\r") if isinstance(line, str) else str(line) for line in before_lines]
    after = [line.rstrip("\n\r") if isinstance(line, str) else str(line) for line in after_lines]

    matcher = difflib.SequenceMatcher(None, before, after)
    changed_before = set()
    changed_after = set()

    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == "equal":
            continue
        if tag in ("replace", "delete"):
            changed_before.update(range(i1, i2))
        if tag in ("replace", "insert"):
            changed_after.update(range(j1, j2))

    result = {
        "before": before,
        "after": after,
        "changed_before_lines": sorted(changed_before),
        "changed_after_lines": sorted(changed_after),
        "summary": {
            "lines_removed": len(changed_before),
            "lines_added": len(changed_after),
            "total_changes": len(changed_before) + len(changed_after),
        },
        "preview_type": preview_type,
    }
    if document_type:
        result["document_type"] = document_type
    if extra:
        result.update(extra)
    return result


def _word_to_lines(file_path):
    doc = Document(file_path)
    lines = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return lines


def compare_word(old_file, new_file):
    return list(
        difflib.unified_diff(
            _word_to_lines(old_file),
            _word_to_lines(new_file),
            fromfile="Baseline Version",
            tofile="Current Version",
            lineterm="",
        )
    )


def format_word_change(old_file, new_file):
    before = _word_to_lines(old_file)
    after = _word_to_lines(new_file)
    return _format_line_comparison(before, after, "document", "word")


def _excel_cell_changes(old_book, new_book):
    changes = []
    all_sheets = set(old_book.sheetnames) | set(new_book.sheetnames)

    for sheet_name in sorted(all_sheets):
        if sheet_name not in old_book.sheetnames:
            changes.append({
                "sheet": sheet_name,
                "cell": "—",
                "before": "(sheet did not exist)",
                "after": "(new sheet added)",
            })
            continue
        if sheet_name not in new_book.sheetnames:
            changes.append({
                "sheet": sheet_name,
                "cell": "—",
                "before": "(sheet existed)",
                "after": "(sheet removed)",
            })
            continue

        old_sheet = old_book[sheet_name]
        new_sheet = new_book[sheet_name]
        rows = max(old_sheet.max_row or 1, new_sheet.max_row or 1)
        cols = max(old_sheet.max_column or 1, new_sheet.max_column or 1)

        for row in range(1, rows + 1):
            for col in range(1, cols + 1):
                old_value = old_sheet.cell(row, col).value
                new_value = new_sheet.cell(row, col).value
                if old_value != new_value:
                    cell = old_sheet.cell(row, col).coordinate
                    changes.append({
                        "sheet": sheet_name,
                        "cell": cell,
                        "before": "" if old_value is None else str(old_value),
                        "after": "" if new_value is None else str(new_value),
                    })
    return changes


def _excel_to_lines(workbook):
    lines = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        lines.append(f"── Sheet: {sheet_name} ──")
        max_row = sheet.max_row or 1
        max_col = sheet.max_column or 1
        for row in range(1, max_row + 1):
            for col in range(1, max_col + 1):
                value = sheet.cell(row, col).value
                if value is not None and str(value).strip():
                    coord = sheet.cell(row, col).coordinate
                    lines.append(f"{coord}: {value}")
        if max_row == 1 and max_col == 1 and sheet.cell(1, 1).value in (None, ""):
            lines.append("(empty sheet)")
    return lines


def compare_excel(old_file, new_file):
    old_book = load_workbook(old_file, data_only=True)
    new_book = load_workbook(new_file, data_only=True)
    return [
        f"{item['sheet']} | {item['cell']} | {item['before']} → {item['after']}"
        for item in _excel_cell_changes(old_book, new_book)
    ]


def format_excel_change(old_file, new_file):
    old_book = load_workbook(old_file, data_only=True)
    new_book = load_workbook(new_file, data_only=True)
    cell_changes = _excel_cell_changes(old_book, new_book)
    result = _format_line_comparison(
        _excel_to_lines(old_book),
        _excel_to_lines(new_book),
        "spreadsheet",
        "excel",
        {"cell_changes": cell_changes},
    )
    result["summary"] = {
        "lines_removed": len(cell_changes),
        "lines_added": len(cell_changes),
        "total_changes": len(cell_changes),
        "cells_changed": len(cell_changes),
    }
    return result


def extract_slide_text(slide):
    text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text.append(shape.text.strip())
    return text


def _ppt_to_lines(file_path):
    presentation = Presentation(file_path)
    lines = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = extract_slide_text(slide)
        if slide_text:
            lines.append(f"── Slide {index} ──")
            lines.extend(slide_text)
    return lines


def compare_powerpoint(old_file, new_file):
    differences = []
    old_ppt = Presentation(old_file)
    new_ppt = Presentation(new_file)
    total = max(len(old_ppt.slides), len(new_ppt.slides))

    for index in range(total):
        if index >= len(old_ppt.slides):
            differences.append(f"[NEW SLIDE] {index + 1}")
            continue
        if index >= len(new_ppt.slides):
            differences.append(f"[DELETED SLIDE] {index + 1}")
            continue

        diff = list(
            difflib.unified_diff(
                extract_slide_text(old_ppt.slides[index]),
                extract_slide_text(new_ppt.slides[index]),
                fromfile=f"Slide {index + 1}",
                tofile=f"Slide {index + 1}",
                lineterm="",
            )
        )
        differences.extend(diff)
    return differences


def format_ppt_change(old_file, new_file):
    before = _ppt_to_lines(old_file)
    after = _ppt_to_lines(new_file)
    return _format_line_comparison(before, after, "presentation", "powerpoint")


def _pdf_to_lines(file_path):
    if PdfReader is None:
        return ["(Install pypdf to preview PDF text)"]
    lines = []
    reader = PdfReader(file_path)
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            lines.append(f"── Page {index} ──")
            lines.extend(line for line in text.splitlines() if line.strip())
    if not lines:
        lines.append("(No extractable text in this PDF — it may be image-based or scanned)")
    return lines


def format_pdf_change(old_file, new_file):
    before = _pdf_to_lines(old_file)
    after = _pdf_to_lines(new_file)
    return _format_line_comparison(before, after, "document", "pdf")


def _rtf_to_lines(file_path):
    if rtf_to_text is None:
        return ["(Install striprtf to preview RTF content)"]
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
            raw = handle.read()
        text = rtf_to_text(raw)
        return [line for line in text.splitlines() if line.strip()]
    except Exception:
        return []


def format_rtf_change(old_file, new_file):
    before = _rtf_to_lines(old_file)
    after = _rtf_to_lines(new_file)
    return _format_line_comparison(before, after, "document", "rtf")


def _xls_to_lines(file_path):
    if xlrd is None:
        return ["(Install xlrd to preview legacy .xls files)"]
    lines = []
    book = xlrd.open_workbook(file_path)
    for sheet in book.sheets():
        lines.append(f"── Sheet: {sheet.name} ──")
        for row_index in range(sheet.nrows):
            for col_index in range(sheet.ncols):
                value = sheet.cell_value(row_index, col_index)
                if value not in (None, ""):
                    col_letter = xlrd.colname(col_index)
                    lines.append(f"{col_letter}{row_index + 1}: {value}")
    return lines


def format_xls_change(old_file, new_file):
    before = _xls_to_lines(old_file)
    after = _xls_to_lines(new_file)
    return _format_line_comparison(before, after, "spreadsheet", "excel")


def format_office_change(old_file, new_file, extension):
    if extension == ".docx":
        return format_word_change(old_file, new_file)
    if extension == ".xlsx":
        return format_excel_change(old_file, new_file)
    if extension == ".pptx":
        return format_ppt_change(old_file, new_file)
    if extension == ".pdf":
        return format_pdf_change(old_file, new_file)
    if extension == ".rtf":
        return format_rtf_change(old_file, new_file)
    if extension == ".xls":
        return format_xls_change(old_file, new_file)
    return None


def extract_office_text(file_path, extension):
    """Readable text preview for document-like files."""
    if extension == ".docx":
        return "\n".join(_word_to_lines(file_path))
    if extension == ".xlsx":
        return "\n".join(_excel_to_lines(load_workbook(file_path, data_only=True)))
    if extension == ".pptx":
        return "\n".join(_ppt_to_lines(file_path))
    if extension == ".pdf":
        return "\n".join(_pdf_to_lines(file_path))
    if extension == ".rtf":
        return "\n".join(_rtf_to_lines(file_path))
    if extension == ".xls":
        return "\n".join(_xls_to_lines(file_path))
    return None
